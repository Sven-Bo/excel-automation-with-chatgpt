import requests
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches

def create_presentation():
    prs = Presentation()
    return prs

def add_title_slide(prs, title, bullet_points, image_url):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = title
    tf = body_shape.text_frame
    for point in bullet_points:
        tf.add_paragraph().text = point

    # download image and add it to the slide
    image_data = requests.get(image_url).content
    with open('dog.jpg', 'wb') as f:
        f.write(image_data)
    left = Inches(5.5)
    top = Inches(2)
    height = Inches(1.5)
    pic = slide.shapes.add_picture('dog.jpg', left, top, height=height)

def main():
    prs = create_presentation()
    bullet_points = ['First bullet point', 'Second bullet point', 'Third bullet point']
    image_url = 'https://place.dog/300/300'
    add_title_slide(prs, 'Introduction', bullet_points, image_url)
    add_title_slide(prs, 'Methods', bullet_points, image_url)
    add_title_slide(prs, 'Results', bullet_points, image_url)

    prs.save('my-awesome-ppt.pptx')

if __name__ == '__main__':
    main()
